from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── 색상 팔레트 ──────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2F, 0x5F)
BLUE   = RGBColor(0x1F, 0x61, 0x8D)
TEAL   = RGBColor(0x02, 0x80, 0x90)
MINT   = RGBColor(0x02, 0xC3, 0x9A)
LTEAL  = RGBColor(0xA8, 0xD8, 0xEA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
DGRAY  = RGBColor(0xD0, 0xD8, 0xE4)
MGRAY  = RGBColor(0x62, 0x65, 0x67)
DARK   = RGBColor(0x1C, 0x28, 0x33)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED    = RGBColor(0xCB, 0x42, 0x35)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN  = RGBColor(0x15, 0x80, 0x3D)

KO = "Malgun Gothic"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = 13.33, 7.5

def I(x): return Inches(x)

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=0.5):
    sp = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line_color:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    return sp

def add_oval(slide, x, y, w, h, fill):
    sp = slide.shapes.add_shape(9, I(x), I(y), I(w), I(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    return sp

def add_text(slide, text, x, y, w, h, size=12, color=None, bold=False,
             align=PP_ALIGN.LEFT, italic=False, valign=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = KO
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = color
    return tb

def section_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, 1.1, fill=NAVY)
    add_rect(slide, 0, 0, 0.14, 1.1, fill=TEAL)
    add_text(slide, title, 0.35, 0.08, 11, 0.62, size=26, color=WHITE, bold=True,
             valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.7, 11, 0.35, size=11, color=LTEAL,
                 valign=MSO_ANCHOR.TOP)

# ─────────────────────────────────────────────────────────────
# S1 표지
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=NAVY)
add_rect(slide, 0, 0, W, 0.12, fill=TEAL)
add_oval(slide, 8.5, -1.2, 5.5, 5.5, TEAL)
add_oval(slide, 9.5, -0.2, 4.5, 4.5, MINT)
add_oval(slide, 10.2,  0.5, 3.0, 3.0, NAVY)
add_rect(slide, 0.7, 1.8, 0.09, 3.2, fill=TEAL)
add_rect(slide, 0.95, 1.9, 3.2, 0.4, fill=TEAL)
add_text(slide, "사업계획서  2026  |  v12 최신버전",
         0.95, 1.9, 3.2, 0.4, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "DentureCare",
         0.95, 2.48, 9, 0.95, size=44, color=WHITE, bold=True)
add_text(slide, "틀니 유지관리 모바일 프로그램",
         0.95, 3.42, 9, 0.6, size=22, color=LTEAL)
add_text(slide, "행동변화 이론 기반 노인 틀니 관리 자동화 시스템",
         0.95, 4.05, 9, 0.48, size=14, color=RGBColor(0xBD,0xD7,0xEA))
add_text(slide, "프로토타입 v12  ·  총 12개 누적 버전  ·  5,000줄 구현",
         0.95, 4.58, 9, 0.42, size=12, color=RGBColor(0x7F,0xB3,0xCC))
add_rect(slide, 0, 6.42, W, 1.08, fill=RGBColor(0x0E,0x1A,0x35))
add_text(slide, "작성: 심미현 (신성대학교 치위생학과 · 가디언즈치과 진료팀장)   |   작성일: 2026년 6월",
         0.7, 6.55, 12, 0.42, size=12, color=LTEAL)

# ─────────────────────────────────────────────────────────────
# S2 목차
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
add_rect(slide, 0, 0, 3.6, H, fill=NAVY)
add_rect(slide, 0, 0, 0.14, H, fill=TEAL)
add_text(slide, "TABLE OF\nCONTENTS", 0.25, 1.6, 3.1, 1.4,
         size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "목  차", 0.25, 3.05, 3.1, 0.85,
         size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, 0.7, 3.92, 2.2, 0.06, fill=TEAL)

toc = [
    ("01", "사업 배경 및 필요성"),
    ("02", "이론적 근거"),
    ("03", "프로그램 개요  (v1 → v12)"),
    ("04", "핵심 기능  (v3~v8)"),
    ("05", "v9~v12 최신 기능"),
    ("06", "4종 액터 시스템"),
    ("07", "리콜 계산 엔진"),
    ("08", "추진 일정 · 예산 · 기대 효과"),
]
for i, (num, title) in enumerate(toc):
    y = 0.5 + i * 0.83
    c = TEAL if i % 2 == 0 else MINT
    # 번호 원
    add_oval(slide, 3.78, y + 0.12, 0.52, 0.52, c)
    add_text(slide, num, 3.78, y + 0.12, 0.52, 0.52,
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 제목 텍스트
    tb = slide.shapes.add_textbox(I(4.45), I(y + 0.08), I(8.5), I(0.48))
    tb.word_wrap = False
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = KO
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = DARK
    r.font.underline = True  # 텍스트 밑줄

# ─────────────────────────────────────────────────────────────
# S3 사업 배경 및 필요성
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "사업 배경 및 필요성", "고령화 사회와 틀니 관리의 중요성")

stats = [
    ("65세 이상\n인구 비율", "19.2%", "2024년 기준 초고령사회 진입 목전", TEAL),
    ("의치 사용\n노인 비율", "31%",   "노인 3명 중 1명이 틀니 사용",       MINT),
    ("정기 검진\n미수검률",  "68%",   "대부분 적절한 리콜 관리 미흡",       YELLOW),
]
for i, (label, stat, desc, clr) in enumerate(stats):
    cx = 0.45 + i * 4.3
    add_rect(slide, cx, 1.28, 3.95, 2.15, fill=LGRAY)
    add_rect(slide, cx, 1.28, 3.95, 0.1, fill=clr)
    add_text(slide, label, cx+0.2, 1.42, 3.55, 0.65, size=13, color=MGRAY)
    add_text(slide, stat,  cx+0.2, 2.0,  3.55, 0.78,
             size=36, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, desc,  cx+0.2, 2.72, 3.55, 0.55, size=11, color=MGRAY,
             align=PP_ALIGN.CENTER)

add_text(slide, "핵심 문제점", 0.45, 3.65, 5.5, 0.42, size=14, color=NAVY, bold=True)
add_rect(slide, 0.45, 4.1, 0.07, 2.8, fill=TEAL)
problems = [
    ("🦷", "관리 지식 부족",  "노인의 68%가 틀니 세척 방법·보관법을 잘못 알고 있음"),
    ("📅", "검진 일정 미준수","리콜 일정을 알지 못해 적절한 검진 시기를 놓침"),
    ("📲", "디지털 지원 부재","행동변화를 유도하는 맞춤형 모바일 도구 전무"),
]
for i, (icon, ttl, desc) in enumerate(problems):
    y = 4.1 + i * 0.95
    add_text(slide, icon + "  " + ttl, 0.65, y,      5.8, 0.42, size=13, color=DARK, bold=True)
    add_text(slide, desc,              0.65, y+0.42, 5.8, 0.42, size=11, color=MGRAY)

add_rect(slide, 7.1, 3.5, 5.9, 3.65, fill=NAVY)
add_rect(slide, 7.1, 3.5, 0.12, 3.65, fill=TEAL)
add_text(slide, "DentureCare의 해결책", 7.38, 3.62, 5.5, 0.48,
         size=14, color=TEAL, bold=True)
needs = [
    "✔ 논문 기반 자동 리콜 일정 계산",
    "✔ 일일 행동 알람으로 관리 습관 형성",
    "✔ 검진 D-day 시각화 및 예약 연결",
    "✔ 응급·미수행 통합 알람 시스템",
    "✔ 건강보험 정보 일일 팝업 안내",
    "✔ 노인 맞춤 글씨 크기·루틴 시간 편집",
]
for i, n in enumerate(needs):
    add_text(slide, n, 7.38, 4.18 + i * 0.49, 5.5, 0.44, size=12, color=WHITE)

# ─────────────────────────────────────────────────────────────
# S4 이론적 근거
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "이론적 근거", "BCW·COM-B 행동변화 이론 + 임상 연구")

add_text(slide, "행동변화 이론 적용", 0.4, 1.2, 6.3, 0.42, size=14, color=NAVY, bold=True)

theories = [
    {
        "color": TEAL,
        "title": "COM-B 모델",
        "meta":  "Susan Michie 교수팀 (UCL) · 2011년 · NHS 공식 채택",
        "concept": "행동(Behaviour)은 능력·기회·동기 세 요소가 상호작용할 때 발생한다는 이론",
        "detail": [
            "능력(Capability): 틀니 세척 동작 수행 + 올바른 관리 지식",
            "기회(Opportunity): 세척 도구 접근성 + 가족·요양보호사 지원",
            "동기(Motivation): 습관(자동적) + 의식적 결정(반성적)",
        ],
    },
    {
        "color": MINT,
        "title": "BCW (Behaviour Change Wheel)",
        "meta":  "Michie 교수팀 · COM-B 확장 프레임워크",
        "concept": "COM-B로 원인 진단 후 최적 중재 전략을 선택하는 행동변화 설계 도구",
        "detail": [
            "교육: 건강보험 정보 팝업 15종으로 지식 제공",
            "환경 재설계: 일일 알람 자동화로 행동 유도",
            "가능화: 1탭 치과 예약 연결로 접근성 향상",
        ],
    },
    {
        "color": BLUE,
        "title": "SRHI (Self-Report Habit Index)",
        "meta":  "Verplanken & Orbell · 2003년 · 12문항 습관 측정 척도",
        "concept": "특정 행동이 얼마나 자동화되었는지(자동성)를 점수화하는 표준 습관 측정 도구",
        "detail": [
            "12개 문항으로 행동 자동성 0~100점 측정",
            "구강보건·운동·식이 습관 연구에서 광범위 활용",
            "DentureCare: 자동성 점수 → 노인 친화적 '관리 점수'로 변환",
        ],
    },
]
CARD_H = 1.78
for i, th in enumerate(theories):
    y = 1.7 + i * (CARD_H + 0.08)
    add_rect(slide, 0.4, y, 6.35, CARD_H, fill=LGRAY)
    add_rect(slide, 0.4, y, 0.14, CARD_H, fill=th["color"])
    # 제목
    add_text(slide, th["title"], 0.68, y+0.06, 5.9, 0.38,
             size=13, color=th["color"], bold=True, valign=MSO_ANCHOR.MIDDLE)
    # 출처
    add_text(slide, th["meta"], 0.68, y+0.44, 5.9, 0.28,
             size=9, color=MGRAY, italic=True, valign=MSO_ANCHOR.TOP)
    # 핵심 개념
    add_text(slide, th["concept"], 0.68, y+0.7, 5.9, 0.34,
             size=10, color=DARK, valign=MSO_ANCHOR.TOP)
    # 세부 항목
    for j, d in enumerate(th["detail"]):
        add_text(slide, "• " + d, 0.72, y+1.03+j*0.24, 5.85, 0.23,
                 size=9.5, color=DARK, valign=MSO_ANCHOR.TOP)

# 오른쪽: 임상 연구
add_rect(slide, 7.0, 1.2, 6.0, 6.0, fill=LGRAY)
add_rect(slide, 7.0, 1.2, 0.12, 6.0, fill=TEAL)
add_text(slide, "임상 연구 근거 (리콜 엔진)", 7.28, 1.28, 5.6, 0.44,
         size=14, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)

papers = [
    ("Tallgren A. (1972)",    "의치 사용자 25년 종단연구\n하악 평균 6.6mm 골 흡수 (상악의 4배)"),
    ("Atwood DA. (1971)",     "잔존치조제흡수(RRR)를 주요 구강질환으로 명명\n발치 후 첫 수개월 흡수 최대"),
    ("Sadr K. et al. (2011)", "의치 장착 후 첫 방문 85.8%에서 점막 손상\n조정 필요 → 적응기 집중 관찰 근거"),
    ("ACP/ADA 가이드라인",    "리라이닝 1~2년마다 권장\n5년 이상 사용 시 연조직 병변 위험 증가"),
]
for i, (ref, content) in enumerate(papers):
    y = 1.88 + i * 1.26
    add_rect(slide, 7.25, y, 5.55, 1.1, fill=WHITE)
    add_rect(slide, 7.25, y, 5.55, 0.05, fill=TEAL)
    add_text(slide, ref,     7.4, y+0.1,  5.2, 0.36, size=11, color=TEAL, bold=True)
    add_text(slide, content, 7.4, y+0.46, 5.2, 0.56, size=10, color=DARK)

# ─────────────────────────────────────────────────────────────
# S5 프로그램 개요
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "프로그램 개요", "DentureCare v12 — 12개 버전, 5,000줄 구현")

add_rect(slide, 0.4, 1.28, 12.55, 1.05, fill=LGRAY)
add_rect(slide, 0.4, 1.28, 0.12, 1.05, fill=TEAL)
add_text(slide,
    "DentureCare는 노인 틀니 사용자의 일일 유지관리 행동을 자동화하고 정기 치과 검진을 유도하는 모바일 헬스 앱으로,\n"
    "BCW·COM-B 이론과 Tallgren·Atwood 임상 연구를 기반으로 설계되었으며 현재 프로토타입 v12까지 구현되었습니다.",
    0.65, 1.32, 12.1, 0.94, size=12, color=DARK)

values = [
    ("자동 일정 계산", "틀니 제작시기 기반\n검진 일정 자동 산출", TEAL),
    ("행동 자동화",   "일일 양치·보관 알람으로\n자기관리 습관 형성",  MINT),
    ("통합 알람",     "응급 자가보고 + 미수행\n자동 감지 통합 시스템", BLUE),
    ("정보 제공",     "건강보험 안내 15종\n하루 한 번 자동 팝업",     NAVY),
]
for i, (ttl, desc, clr) in enumerate(values):
    cx = 0.4 + i * 3.25
    add_rect(slide, cx, 2.5, 3.05, 2.05, fill=clr)
    add_text(slide, ttl,  cx+0.1, 2.6,  2.85, 0.52,
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, cx+0.3, 3.1, 2.45, 0.04, fill=WHITE)
    add_text(slide, desc, cx+0.1, 3.2,  2.85, 0.78,
             size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "버전별 발전 이력", 0.4, 4.75, 12, 0.38, size=13, color=NAVY, bold=True)

versions = [
    ("v1~2", "기초알람\nSOS",    MGRAY),
    ("v3",   "등록+\n리콜",      BLUE),
    ("v4",   "영구\n저장",       BLUE),
    ("v5",   "단골치과\n검색",   TEAL),
    ("v6~7", "검진탭\n재구성",   TEAL),
    ("v8",   "데이터\n영구저장", TEAL),
    ("v9",   "팝업\n알람",       MINT),
    ("v10",  "통합\n알람",       MINT),
    ("v11",  "수정\n다이얼로그", MINT),
    ("v12",  "일일정보\n팝업+UI",MINT),
]
n = len(versions)
gap = 12.55 / n
for i, (v, desc, clr) in enumerate(versions):
    x = 0.38 + i * gap
    add_oval(slide, x, 5.25, 0.75, 0.75, clr)
    add_text(slide, v, x, 5.25, 0.75, 0.75,
             size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x-0.08, 6.06, 0.91, 0.56, size=8, color=DARK,
             align=PP_ALIGN.CENTER)
    if i < n-1:
        add_rect(slide, x+0.75, 5.59, gap-0.75, 0.05, fill=DGRAY)

add_text(slide, "◀ 현재", 0.38 + 9*gap + 0.78, 5.28, 1.2, 0.34,
         size=9.5, color=RED, bold=True)

# ─────────────────────────────────────────────────────────────
# S6 핵심 기능 (v3~v8)
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "핵심 기능 (v3~v8)", "환자 등록 → 리콜 계산 → 검진 유도 → 데이터 저장")

features = [
    {"num":"01","title":"환자 등록 시스템  (v3)","color":TEAL,
     "items":["6단계 등록 마법사 (이름·출생연도·틀니제작시기)",
              "보호자 초대 코드 (카카오톡 공유)",
              "단골 치과 등록 및 검색 (v5)",
              "미입력 항목 자동 안내 (v6)"]},
    {"num":"02","title":"리콜 자동 계산 엔진  (v5)","color":MINT,
     "items":["0~3개월 적응기: 1개월 간격 (집중 관찰)",
              "3~12개월 1년차: 3개월 간격 (급격 흡수기)",
              "1~5년 안정기: 6개월 간격",
              "5년+ 장기사용: 6개월 간격 (교체 검토)"]},
    {"num":"03","title":"검진 유도 시스템  (v7)","color":BLUE,
     "items":["검진 타임라인 시각화 (D-day 바 그래프)",
              "D-14 노랑 / D-7 주황 / D-3 빨강 강조",
              "가디언즈치과 1탭 예약 연결",
              "검진 점검 항목: 적합도·잇몸·점막·리라이닝"]},
    {"num":"04","title":"데이터 영구 저장  (v8)","color":NAVY,
     "items":["환자 정보·틀니제작시기·단골치과 영구 보존",
              "연속일수·관리점수 자동 누적",
              "앱 재실행 시 자동 불러오기",
              "데이터 초기화 메뉴 (확인 다이얼로그 포함)"]},
]
for i, feat in enumerate(features):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.58
    y = 1.22 + row * 2.97
    add_rect(slide, x, y, 6.28, 2.78, fill=LGRAY)
    add_rect(slide, x, y, 6.28, 0.09, fill=feat["color"])
    add_oval(slide, x+0.18, y+0.2, 0.52, 0.52, feat["color"])
    add_text(slide, feat["num"], x+0.18, y+0.2, 0.52, 0.52,
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, feat["title"], x+0.85, y+0.22, 5.25, 0.44,
             size=13, color=feat["color"], bold=True, valign=MSO_ANCHOR.MIDDLE)
    for j, item in enumerate(feat["items"]):
        add_text(slide, "• "+item, x+0.28, y+0.84+j*0.46, 5.82, 0.4,
                 size=11, color=DARK, valign=MSO_ANCHOR.TOP)

# ─────────────────────────────────────────────────────────────
# S7 v9~v12 최신 기능
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "v9~v12 최신 기능", "2026년 상반기 추가 구현 — 4개 버전, 핵심 UX 고도화")

new_feats = [
    {"ver":"v9",  "title":"리콜 임박 팝업 알람",          "color":YELLOW,
     "items":["D-14: 노란색 팝업 '2주 뒤 검진 예정이에요'",
              "D-7: 주황색 + 예약 버튼 '이번 주 검진이에요!'",
              "D-3: 빨간색 + 상단 배너 '3일 뒤 검진이에요!'",
              "D-day: 최고 강조 / D+1 이후 '다녀오셨어요?' 팝업",
              "하루 한 번만 표시 (알림 피로 방지)"]},
    {"ver":"v10", "title":"통합 알람 시스템 (응급 + 미수행)","color":RED,
     "items":["자가 보고 메뉴: 잇몸통증·틀니파손·틀니헐거움",
              "미수행 자동 감지: 1일·3일 미수행 단계별 알람",
              "취침 전 틀니 탈거 확인 알람",
              "응급별 맞춤 대처법·주의사항 안내",
              "1탭으로 가디언즈치과 바로 연결"]},
    {"ver":"v11", "title":"틀니 제작시기 수정 다이얼로그",  "color":PURPLE,
     "items":["설정 화면에서 틀니 제작 연도·월 수정 가능",
              "수정 시 확인 다이얼로그 ('중요' 뱃지 표시)",
              "수정 후 리콜 일정 자동 재계산 및 알림",
              "입력 오류 방지용 범위 검증 로직"]},
    {"ver":"v12", "title":"일일 정보 팝업 + 노인 맞춤 UI", "color":MINT,
     "items":["건강보험 유지관리 안내 7종 (만 65세 이상)",
              "틀니 재제작 조건 안내 3종",
              "일반 안내 5종 (처음 3개월 무료 점검 등)",
              "글씨 크기 조절: 소·중·대 3단계 (v12.1)",
              "루틴 알람 시간 편집 + 시·분 선택 팝업 (v12.2)"]},
]
for i, feat in enumerate(new_feats):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.58
    y = 1.22 + row * 2.97
    add_rect(slide, x, y, 6.28, 2.78, fill=LGRAY)
    add_rect(slide, x, y, 6.28, 0.09, fill=feat["color"])
    add_rect(slide, x+0.18, y+0.18, 0.72, 0.4, fill=feat["color"])
    add_text(slide, feat["ver"], x+0.18, y+0.18, 0.72, 0.4,
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, feat["title"], x+1.02, y+0.2, 5.1, 0.44,
             size=13, color=feat["color"], bold=True, valign=MSO_ANCHOR.MIDDLE)
    for j, item in enumerate(feat["items"]):
        add_text(slide, "• "+item, x+0.28, y+0.76+j*0.4, 5.82, 0.37,
                 size=10.5, color=DARK, valign=MSO_ANCHOR.TOP)

# ─────────────────────────────────────────────────────────────
# S8 4종 액터
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "4종 액터 시스템", "노인·가족·요양보호사·관리자의 역할별 인터페이스")

actors = [
    {"code":"A1","name":"노인 본인",  "color":TEAL,
     "roles":["일일 행동 체크 (양치·보관)","검진 D-day 확인 및 예약",
              "응급 자가 보고 (v10)","관리 점수·연속일수 확인"],
     "screens":"홈·일일관리·검진·기록 탭"},
    {"code":"A2","name":"가족",       "color":MINT,
     "roles":["어르신 진행률 원격 확인","격려 메시지 발송",
              "응급 알림 수신","검진 일정 공유 확인"],
     "screens":"가족 대시보드·알림"},
    {"code":"A3","name":"요양보호사", "color":BLUE,
     "roles":["다수 어르신 동시 관리","시설 응답률 통계",
              "검진 미수검자 알림","일괄 체크인 기능"],
     "screens":"담당자 리스트·시설 통계"},
    {"code":"A5","name":"시스템 관리자","color":NAVY,
     "roles":["전국 사용자·DAU 현황","47개 시설 랭킹 관리",
              "이상 패턴 감지","KPI·CSV 익스포트"],
     "screens":"대시보드·시설별 성과·설정"},
]
for i, actor in enumerate(actors):
    x = 0.38 + i * 3.27
    add_rect(slide, x, 1.28, 3.04, 5.9, fill=LGRAY)
    add_rect(slide, x, 1.28, 3.04, 1.05, fill=actor["color"])
    add_oval(slide, x+1.16, 1.33, 0.72, 0.72, WHITE)
    add_text(slide, actor["code"], x+1.16, 1.33, 0.72, 0.72,
             size=13, color=actor["color"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, actor["name"], x+0.1, 2.1, 2.84, 0.46,
             size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "주요 역할", x+0.22, 2.52, 2.6, 0.36,
             size=11, color=actor["color"], bold=True)
    for j, role in enumerate(actor["roles"]):
        add_text(slide, "▸ "+role, x+0.22, 2.9+j*0.54, 2.72, 0.46,
                 size=10.5, color=DARK, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, x+0.18, 5.12, 2.68, 0.03, fill=DGRAY)
    add_text(slide, "화면 구성", x+0.22, 5.2, 2.6, 0.3, size=10, color=MGRAY, bold=True)
    add_text(slide, actor["screens"], x+0.22, 5.5, 2.72, 0.55, size=10.5, color=DARK)

# ─────────────────────────────────────────────────────────────
# S9 리콜 계산 엔진
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "리콜 계산 엔진", "논문 근거 기반 4단계 자동 검진 주기 + v9 임박 팝업")

stages = [
    ("0~3개월",  "적응기\n(집중 관찰)",  "1개월", "점막 적응 시기\n85.8% 조정 필요",  TEAL),
    ("3~12개월", "1년차\n(급격 흡수기)", "3개월", "골 흡수 최대\n적합도 변화 시작",  MINT),
    ("1~5년",    "안정기",              "6개월", "골 흡수 안정\n리라이닝 검토",       BLUE),
    ("5년 이상", "장기사용\n(교체 검토)","6개월", "부적합·병변 위험↑\n교체 강력 권고", NAVY),
]
for i, (period, stage, interval, desc, clr) in enumerate(stages):
    x = 0.38 + i * 3.27
    add_rect(slide, x, 1.28, 3.04, 3.55, fill=clr)
    add_text(slide, period,   x+0.1, 1.35, 2.84, 0.42,
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x+0.3, 1.77, 2.44, 0.04, fill=WHITE)
    add_text(slide, stage,    x+0.1, 1.84, 2.84, 0.72,
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, x+0.38, 2.62, 2.28, 0.84, fill=WHITE)
    add_text(slide, "리콜 주기", x+0.38, 2.65, 2.28, 0.3,
             size=10, color=clr, align=PP_ALIGN.CENTER)
    add_text(slide, interval,  x+0.38, 2.92, 2.28, 0.44,
             size=22, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, desc,      x+0.15, 3.58, 2.74, 0.62,
             size=10.5, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        add_rect(slide, x+3.04, 2.96, 0.23, 0.08, fill=DGRAY)

# D-day 팝업 섹션
add_rect(slide, 0.38, 5.02, 12.57, 2.2, fill=LGRAY)
add_rect(slide, 0.38, 5.02, 0.12, 2.2, fill=YELLOW)
add_text(slide, "v9 리콜 임박 팝업 알람  —  D-day 단계별 강도",
         0.65, 5.1, 7, 0.4, size=13, color=NAVY, bold=True)

dday = [
    ("D-14 이내",  "노란색\n'2주 뒤 검진 예정'",   YELLOW),
    ("D-7 이내",   "주황색\n+ 예약 버튼",           ORANGE),
    ("D-3 이내",   "빨간색\n+ 상단 배너",           RED),
    ("D-day 당일", "최고 강조\n+ 팝업 알람",        RGBColor(0x96,0x28,0x1C)),
    ("D+1 이후",   "'다녀오셨어요?'\n재설정 팝업",  PURPLE),
]
for i, (label, action, clr) in enumerate(dday):
    x = 0.52 + i * 2.52
    add_rect(slide, x, 5.6, 2.35, 1.0, fill=clr)
    add_text(slide, label,  x+0.12, 5.65, 2.11, 0.38,
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, action, x+0.12, 6.02, 2.11, 0.48,
             size=9.5, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────
# S10 추진 일정
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "추진 일정", "v12까지 완료 · v13+ 차기 개발 계획")

phases = [
    ("Phase 1\n기반 구축",    "2026년 1~2월", MGRAY, "완료",
     ["요구사항 분석 및 이론 검토","v1~v2 기초 알람·SOS 구현","BCW·COM-B 적용 설계"],
     ["v1·v2 프로토타입","이론 적용 설계서"]),
    ("Phase 2\nv3~v8 핵심",  "2026년 3~4월", TEAL,  "완료",
     ["환자 등록·리콜 계산 엔진","검진 탭 개편 및 검진 유도","데이터 영구 저장 구현"],
     ["v3~v8 앱 소스코드","리콜·저장 모듈"]),
    ("Phase 3\nv9~v12 고도화","2026년 5~6월", MINT,  "완료",
     ["팝업 알람·통합 알람 시스템","건강보험 정보 팝업 15종","글씨 크기·루틴 시간 편집"],
     ["v9~v12 소스코드","건강보험 콘텐츠"]),
    ("Phase 4\nv13+ 출시",    "2026년 하반기",NAVY,  "예정",
     ["치과 전문가 모드 (A4)","EMR 연계 (HL7 FHIR)","앱스토어 정식 출시"],
     ["정식 앱 배포","시설 협약서"]),
]
add_rect(slide, 0.38, 3.85, 12.57, 0.08, fill=DGRAY)
for i, (phase, period, clr, status, tasks, outputs) in enumerate(phases):
    x = 0.38 + i * 3.27
    add_oval(slide, x+1.02, 3.6, 0.62, 0.62, clr)
    add_text(slide, str(i+1), x+1.02, 3.6, 0.62, 0.62,
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 상단 박스
    add_rect(slide, x+0.14, 1.28, 2.9, 2.25, fill=clr)
    badge_clr = GREEN if status == "완료" else ORANGE
    add_rect(slide, x+1.62, 1.35, 1.24, 0.32, fill=badge_clr)
    add_text(slide, "● "+status, x+1.62, 1.35, 1.24, 0.32,
             size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, phase,  x+0.24, 1.7, 2.7, 0.8,
             size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, period, x+0.24, 2.44, 2.7, 0.34,
             size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x+0.38, 2.78, 2.34, 0.04, fill=WHITE)
    for j, task in enumerate(tasks):
        add_text(slide, "• "+task, x+0.24, 2.86+j*0.22, 2.7, 0.2,
                 size=8.5, color=WHITE, valign=MSO_ANCHOR.TOP)
    # 하단 박스
    add_rect(slide, x+0.14, 4.32, 2.9, 2.88, fill=LGRAY)
    add_rect(slide, x+0.14, 4.32, 2.9, 0.07, fill=clr)
    add_text(slide, "주요 산출물", x+0.24, 4.44, 2.7, 0.36,
             size=10.5, color=clr, bold=True)
    for j, out in enumerate(outputs):
        add_text(slide, "▸ "+out, x+0.24, 4.88+j*0.76, 2.7, 0.6,
                 size=10.5, color=DARK, valign=MSO_ANCHOR.TOP)

# ─────────────────────────────────────────────────────────────
# S11 예산 계획
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "예산 계획", "단계별 예산 배분 및 소요 비용")

add_text(slide, "항목별 예산 (단위: 만원)", 0.4, 1.28, 12.6, 0.44,
         size=14, color=NAVY, bold=True)

# 열 정의 (시작 x, 너비) — 합계 = 12.5
COL = [(0.4, 1.55), (1.95, 2.75), (4.7, 6.0), (10.7, 2.05)]
# 항목 | 내용 | 산출 근거 | 금액(만원)
HDR_LABELS = ["항목", "내용", "산출 근거", "금액(만원)"]
HDR_Y = 1.82
ROW_H = 0.6
PAD_X = 0.12  # 셀 내 좌우 여백

# 헤더
add_rect(slide, 0.4, HDR_Y, 12.35, 0.46, fill=NAVY)
for (cx, cw), lbl in zip(COL, HDR_LABELS):
    add_text(slide, lbl, cx+PAD_X, HDR_Y+0.02, cw-PAD_X*2, 0.42,
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

budget = [
    ("개발비",     "앱 개발·UI/UX 디자인",  "프리랜서 개발자 300만원 × 4개월",       "1,200"),
    ("서버·인프라", "클라우드 서버·DB·보안", "AWS 월 30만원 × 12개월",                "360"),
    ("파일럿 운영", "테스트 노인 20명",      "교통비·사례비 10만원 × 20명",           "200"),
    ("홍보·마케팅", "앱스토어 등록·홍보물",  "앱스토어 등록비 + 리플렛·포스터 제작",  "150"),
    ("기타 운영",   "인건비·교육·관리",      "연구자 활동비, 워크숍·교육비 등",       "290"),
]
for i, (cat, desc, basis, amt) in enumerate(budget):
    ry = HDR_Y + 0.46 + i * ROW_H
    bg = WHITE if i % 2 == 0 else LGRAY
    add_rect(slide, 0.4, ry, 12.35, ROW_H, fill=bg)
    # 열 구분선
    for j in range(1, 4):
        add_rect(slide, COL[j][0]-0.02, ry, 0.02, ROW_H, fill=DGRAY)
    # 셀 텍스트
    add_text(slide, cat,   COL[0][0]+PAD_X, ry, COL[0][1]-PAD_X, ROW_H,
             size=11, color=DARK, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, desc,  COL[1][0]+PAD_X, ry, COL[1][1]-PAD_X, ROW_H,
             size=10.5, color=MGRAY, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, basis, COL[2][0]+PAD_X, ry, COL[2][1]-PAD_X, ROW_H,
             size=10.5, color=DARK, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, amt,   COL[3][0]+PAD_X, ry, COL[3][1]-PAD_X*2, ROW_H,
             size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# 합계 행
total_y = HDR_Y + 0.46 + 5 * ROW_H
add_rect(slide, 0.4, total_y, 12.35, 0.5, fill=NAVY)
add_text(slide, "합  계", COL[0][0]+PAD_X, total_y, 9.8, 0.5,
         size=12, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
add_text(slide, "2,200만원", COL[3][0]+PAD_X, total_y, COL[3][1]-PAD_X*2, 0.5,
         size=13, color=MINT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# 하단: 파이 차트 + 예산 확보
chart_data = ChartData()
chart_data.categories = ["개발비", "서버·인프라", "파일럿", "홍보", "기타"]
chart_data.add_series("예산", (1200, 360, 200, 150, 290))
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, I(0.4), I(5.32), I(5.8), I(2.0), chart_data)
co = chart.chart
co.has_legend = True; co.legend.position = 2
co.plots[0].has_data_labels = True
co.plots[0].data_labels.show_percent = True

add_rect(slide, 6.4, 5.32, 6.55, 2.0, fill=LGRAY)
add_rect(slide, 6.4, 5.32, 0.12, 2.0, fill=TEAL)
add_text(slide, "예산 확보 계획", 6.65, 5.4, 6.2, 0.4, size=13, color=NAVY, bold=True)
sources = [
    "• 교내 창업지원금 (신성대 산학처): 700만원",
    "• 복지부 디지털 헬스케어 공모사업: 1,000만원",
    "• 자체 시범 운영 수익: 500만원",
]
for i, s in enumerate(sources):
    add_text(slide, s, 6.65, 5.88+i*0.46, 6.2, 0.42, size=11.5, color=DARK,
             valign=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────
# S12 기대 효과
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=WHITE)
section_header(slide, "기대 효과", "임상적·사회적·경제적 효과")

effects = [
    {"cat":"임상적 효과","color":TEAL,
     "items":[("검진 수검률 향상","리콜 자동 알람으로 정기 검진율 30%↑ 목표"),
              ("구강 건강 개선",  "올바른 틀니 관리 습관으로 점막 병변 감소"),
              ("응급 대응 단축",  "1탭 응급 연결로 대처 시간 50% 단축")]},
    {"cat":"사회적 효과","color":MINT,
     "items":[("노인 자립 지원",   "자기관리 능력 향상으로 의존도 감소"),
              ("가족 부담 경감",   "원격 모니터링으로 보호자 불안 해소"),
              ("요양기관 효율화", "요양보호사의 다중 관리 시스템화")]},
    {"cat":"경제적 효과","color":BLUE,
     "items":[("의료비 절감",     "예방적 관리로 응급 치료비 절감"),
              ("치과 내원 최적화","불필요한 응급 내원 감소, 예약제 전환"),
              ("디지털 헬스 확장","요양시설 연계 확대 가능성")]},
]
for i, eff in enumerate(effects):
    x = 0.38 + i * 4.35
    add_rect(slide, x, 1.28, 4.1, 5.92, fill=LGRAY)
    add_rect(slide, x, 1.28, 4.1, 0.56, fill=eff["color"])
    add_text(slide, eff["cat"], x+0.2, 1.36, 3.7, 0.42,
             size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, (ttl, desc) in enumerate(eff["items"]):
        y = 2.02 + j * 1.66
        add_rect(slide, x+0.2, y, 3.7, 1.46, fill=WHITE)
        add_rect(slide, x+0.2, y, 3.7, 0.07, fill=eff["color"])
        add_text(slide, ttl,  x+0.36, y+0.12, 3.38, 0.4,
                 size=12, color=eff["color"], bold=True)
        add_text(slide, desc, x+0.36, y+0.54, 3.38, 0.78,
                 size=11, color=DARK, valign=MSO_ANCHOR.TOP)

# ─────────────────────────────────────────────────────────────
# S13 결론
# ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, fill=NAVY)
add_rect(slide, 0, 0, W, 0.1, fill=TEAL)
add_rect(slide, 0, H-0.1, W, 0.1, fill=TEAL)
add_oval(slide,  8.0, -0.5, 6.0, 6.0, TEAL)
add_oval(slide,  9.5,  0.5, 5.0, 5.0, MINT)
add_oval(slide, 10.5,  1.5, 3.5, 3.5, BLUE)

add_text(slide, "결  론", 0.8, 0.82, 8, 0.68, size=32, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.56, 3.5, 0.09, fill=TEAL)
add_text(slide,
    "DentureCare는 BCW·COM-B 행동변화 이론과 임상 연구(Tallgren 1972, Atwood 1971)를 기반으로\n"
    "v12까지 12개 버전, 5,000줄 이상 구현된 노인 틀니 유지관리 통합 모바일 시스템입니다.",
    0.8, 1.78, 8.5, 1.0, size=13, color=LTEAL)

achievements = [
    ("v3~v8 핵심 기반",     "등록·리콜 계산·검진 유도·영구 저장 구축"),
    ("v9 팝업 알람",        "D-day 단계별 리콜 임박 자동 팝업"),
    ("v10 통합 알람",       "응급 자가보고 + 미수행 자동 감지"),
    ("v11 수정 다이얼로그", "틀니 제작시기 수정 → 리콜 자동 재계산"),
    ("v12 일일 정보 팝업",  "건강보험 안내 15종 + 글씨 크기·루틴 편집"),
]
add_text(slide, "v12까지 핵심 성취", 0.8, 2.98, 6, 0.4, size=14, color=TEAL, bold=True)
for i, (ttl, desc) in enumerate(achievements):
    y = 3.48 + i * 0.6
    add_oval(slide, 0.8, y+0.08, 0.38, 0.38, TEAL)
    add_text(slide, "✓", 0.8, y+0.08, 0.38, 0.38,
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, ttl+"  —  "+desc, 1.3, y+0.06, 6.6, 0.4, size=11.5, color=WHITE,
             valign=MSO_ANCHOR.MIDDLE)

add_text(slide, "다음 단계: v13 치과 전문가 모드 → EMR 연계 (HL7 FHIR) → 앱스토어 정식 출시",
         0.8, 6.38, 8.5, 0.44, size=12, color=LTEAL)
add_text(slide,
    "\"행동 자동화 + 검진 자동 유도\"\n통합 시스템 완성 (v12)",
    8.5, 4.5, 4.5, 1.5, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ── 저장 ──────────────────────────────────────────────────────
output = "/Users/simmihyeon/Desktop/전공심화/Denture Care program/DentureCare_사업계획서.pptx"
prs.save(output)
print("저장 완료:", output)
print("총 슬라이드:", len(prs.slides))
